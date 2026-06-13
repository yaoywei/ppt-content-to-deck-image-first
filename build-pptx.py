#!/usr/bin/env python3
"""
build-pptx.py — 用 python-pptx 拼 15 张图成 .pptx
- 16:9 (13.333" x 7.5")
- 每页一张图作为 full-bleed background
- 加上 page header (章节编号+标题) 作为 overlay text
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from PIL import Image

EXAMPLES_DIR = os.path.expanduser("~/ppt-content-to-deck-image-first-build/examples")
OUTPUT_DIR = os.path.expanduser("~/ppt-content-to-deck-image-first-build")

# 15 页的元数据 (文件名 + 章节编号 + 标题)
PAGES = [
    ("cover-red-diagonal.png", None, None),  # 封面没有 header
    ("p02-company-overview.png", "01", "公司概况"),
    ("p03-mission-positioning.png", "02", "企业使命与定位"),
    ("p04-qualifications.png", "03", "核心资质与荣誉"),
    ("p05-business-overview.png", "04", "主营业务"),
    ("p06-facade-cleaning.png", "05", "城市场景 - 无人机外墙清洗"),
    ("p07-municipal-inspection.png", "06", "城市场景 - 市政设施巡检"),
    ("p08-drug-eradication.png", "07", "城市场景 - 公安禁毒航测"),
    ("p09-power-inspection.png", "08", "工业场景 - 电力智能巡检"),
    ("p10-petrochemical-inspection.png", "09", "工业场景 - 石油化工巡检"),
    ("p11-international-business.png", "10", "国际低空产业资源对接"),
    ("p12-training-program.png", "11", "鲲鹏标准 - 职业技能培训"),
    ("p13-core-team.png", "12", "核心管理团队"),
    ("p14-history-timeline.png", "13", "发展历程"),
    ("p15-contact-us.png", "14", "联系我们"),
]

# 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 空白 layout (index 6 = blank)
blank_layout = prs.slide_layouts[6]

for filename, section_num, section_title in PAGES:
    slide = prs.slides.add_slide(blank_layout)
    img_path = os.path.join(EXAMPLES_DIR, filename)
    # Full-bleed image
    slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
    # 封面不加 header overlay (图片已含)
    if section_num is None:
        continue

output_path = os.path.join(OUTPUT_DIR, "kunpeng-yihang-v4-leadership-style.pptx")
prs.save(output_path)
print(f"✅ Saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Size: {os.path.getsize(output_path) / 1024 / 1024:.1f} MB")
